"""Generate a PowerPoint overview of the NeqSim agent and skill ecosystem."""

from __future__ import annotations

import argparse
from collections import Counter
from datetime import date
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(18, 35, 52)
INK = RGBColor(29, 43, 55)
MUTED = RGBColor(89, 107, 120)
CYAN = RGBColor(0, 151, 194)
TEAL = RGBColor(0, 128, 122)
CORAL = RGBColor(230, 93, 73)
YELLOW = RGBColor(245, 184, 65)
PALE_BLUE = RGBColor(226, 243, 248)
PALE_TEAL = RGBColor(225, 243, 239)
PALE_CORAL = RGBColor(252, 235, 230)
PALE_YELLOW = RGBColor(254, 245, 217)
OFF_WHITE = RGBColor(247, 249, 250)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(207, 218, 224)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Cascadia Mono"


class RepoCounts:
    """Container for repository inventory counts used in the slides."""

    def __init__(self, workspace_root: Path) -> None:
        """Build counts from checked-out repository files.

        Args:
            workspace_root: Parent directory containing the NeqSim repositories.
        """
        self.workspace_root = workspace_root
        self.core_agents = self._count("neqsim/.github/agents/*.agent.md")
        self.core_skills = self._count("neqsim/.github/skills/*/SKILL.md")
        self.community_agents = self._count(
            "neqsim-community-agents/agents/*/agent.yaml"
        )
        self.enterprise_agents = self._count(
            "neqsim-enterprise-agents/agents/*/agent.yaml"
        )
        self.community_skills = self._count(
            "neqsim-community-skills/skills/*/*/SKILL.md"
        )
        self.enterprise_skills = self._count(
            "neqsim-enterprise-skills/skills/*/*/SKILL.md"
        )
        self.community_skill_categories = self._category_counts(
            workspace_root / "neqsim-community-skills" / "skills"
        )
        self.enterprise_skill_categories = self._category_counts(
            workspace_root / "neqsim-enterprise-skills" / "skills"
        )

    def _count(self, pattern: str) -> int:
        """Count files matching a workspace-relative glob.

        Args:
            pattern: Workspace-relative glob pattern.

        Returns:
            Number of matching files.
        """
        return len(list(self.workspace_root.glob(pattern)))

    @staticmethod
    def _category_counts(skills_root: Path) -> Counter[str]:
        """Count packaged skills by their category directory.

        Args:
            skills_root: Root directory containing category subdirectories.

        Returns:
            Counts keyed by category name.
        """
        counts: Counter[str] = Counter()
        if not skills_root.exists():
            return counts
        for skill_file in skills_root.glob("*/*/SKILL.md"):
            counts[skill_file.parents[1].name] += 1
        return counts

    @property
    def total_agents(self) -> int:
        """Return the number of checked-out agent package definitions."""
        return self.core_agents + self.community_agents + self.enterprise_agents

    @property
    def total_skills(self) -> int:
        """Return the number of checked-out skill package definitions."""
        return self.core_skills + self.community_skills + self.enterprise_skills


class DeckBuilder:
    """Small layout helper for the NeqSim organization presentation."""

    def __init__(self, counts: RepoCounts) -> None:
        """Initialize the presentation and shared theme.

        Args:
            counts: Current repository inventory.
        """
        self.counts = counts
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    def add_slide(self, title: str, kicker: str, number: int):
        """Create a standard content slide.

        Args:
            title: Slide title.
            kicker: Short section label above the title.
            number: Slide number shown in the footer.

        Returns:
            Newly created slide.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._background(slide, OFF_WHITE)
        self.text(slide, 0.62, 0.34, 7.8, 0.28, kicker.upper(), 10, CYAN, bold=True)
        self.text(slide, 0.62, 0.68, 11.8, 0.64, title, 26, NAVY, bold=True)
        self.line(slide, 0.62, 1.38, 12.08, 1.38, LINE, 1.2)
        self.text(slide, 0.62, 7.12, 8.8, 0.18, "NeqSim agent and skill ecosystem", 8, MUTED)
        self.text(slide, 12.0, 7.08, 0.65, 0.22, f"{number:02d}", 9, MUTED, align=PP_ALIGN.RIGHT)
        return slide

    @staticmethod
    def _background(slide, color: RGBColor) -> None:
        """Set a solid slide background.

        Args:
            slide: PowerPoint slide.
            color: RGB fill color.
        """
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def text(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        size: float,
        color: RGBColor = INK,
        bold: bool = False,
        font: str = BODY_FONT,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        margin: float = 0.0,
    ):
        """Add a text box with predictable typography.

        Args:
            slide: PowerPoint slide.
            x: Left position in inches.
            y: Top position in inches.
            w: Width in inches.
            h: Height in inches.
            value: Text content.
            size: Font size in points.
            color: Font color.
            bold: Whether to use bold text.
            font: Font family.
            align: Paragraph alignment.
            valign: Vertical anchor.
            margin: Internal margin in inches.

        Returns:
            Created text box shape.
        """
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.word_wrap = True
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = value
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def rect(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: RGBColor,
        radius: bool = True,
        line_color: RGBColor | None = None,
        line_width: float = 0.8,
    ):
        """Add a rectangle or rounded rectangle.

        Args:
            slide: PowerPoint slide.
            x: Left position in inches.
            y: Top position in inches.
            w: Width in inches.
            h: Height in inches.
            fill: Fill color.
            radius: Whether to use rounded corners.
            line_color: Optional border color.
            line_width: Border width in points.

        Returns:
            Created shape.
        """
        shape_type = (
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
            if radius
            else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        )
        shape = slide.shapes.add_shape(
            shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line_color is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(line_width)
        return shape

    def line(
        self,
        slide,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        color: RGBColor,
        width: float = 1.0,
        arrow: bool = False,
    ):
        """Add a straight connector.

        Args:
            slide: PowerPoint slide.
            x1: Start x in inches.
            y1: Start y in inches.
            x2: End x in inches.
            y2: End y in inches.
            color: Line color.
            width: Line width in points.
            arrow: Whether to add a terminal arrowhead.

        Returns:
            Created connector.
        """
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        connector.line.color.rgb = color
        connector.line.width = Pt(width)
        if arrow:
            connector.line.end_arrowhead = True
        return connector

    def pill(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        label: str,
        fill: RGBColor,
        color: RGBColor = WHITE,
    ) -> None:
        """Add a compact label pill.

        Args:
            slide: PowerPoint slide.
            x: Left position in inches.
            y: Top position in inches.
            w: Width in inches.
            label: Label text.
            fill: Fill color.
            color: Text color.
        """
        self.rect(slide, x, y, w, 0.34, fill)
        self.text(
            slide,
            x,
            y + 0.01,
            w,
            0.27,
            label,
            9,
            color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

    def bullet_list(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        items: Sequence[str],
        size: float = 13,
        color: RGBColor = INK,
        gap: float = 0.46,
    ) -> None:
        """Draw a compact bullet list using colored markers.

        Args:
            slide: PowerPoint slide.
            x: Left position in inches.
            y: Top position in inches.
            w: Width in inches.
            items: Bullet strings.
            size: Font size in points.
            color: Text color.
            gap: Vertical distance between bullets in inches.
        """
        for index, item in enumerate(items):
            item_y = y + index * gap
            self.rect(slide, x, item_y + 0.11, 0.09, 0.09, CYAN, radius=False)
            self.text(slide, x + 0.2, item_y, w - 0.2, 0.36, item, size, color)

    def metric(self, slide, x: float, y: float, value: str, label: str, fill: RGBColor) -> None:
        """Add a large metric tile.

        Args:
            slide: PowerPoint slide.
            x: Left position in inches.
            y: Top position in inches.
            value: Main metric value.
            label: Metric description.
            fill: Tile fill color.
        """
        self.rect(slide, x, y, 2.72, 1.12, fill)
        self.text(slide, x + 0.18, y + 0.12, 2.34, 0.52, value, 27, NAVY, bold=True)
        self.text(slide, x + 0.18, y + 0.68, 2.34, 0.27, label, 10, MUTED)

    def add_title_slide(self) -> None:
        """Add the title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._background(slide, NAVY)
        self.rect(slide, 0.0, 0.0, 0.22, 7.5, CYAN, radius=False)
        self.rect(slide, 10.45, 0.0, 2.88, 7.5, RGBColor(13, 48, 68), radius=False)
        for x, y, color in [
            (10.93, 1.18, CYAN),
            (11.97, 2.28, TEAL),
            (10.75, 3.66, CORAL),
            (12.08, 5.10, YELLOW),
        ]:
            self.rect(slide, x, y, 0.58, 0.58, color)
        self.line(slide, 11.50, 1.47, 11.97, 2.57, WHITE, 1.3)
        self.line(slide, 11.97, 2.86, 11.04, 3.95, WHITE, 1.3)
        self.line(slide, 11.04, 4.24, 12.08, 5.39, WHITE, 1.3)
        self.text(slide, 0.83, 0.78, 8.8, 0.28, "NEQSIM ECOSYSTEM", 11, CYAN, bold=True)
        self.text(slide, 0.83, 1.42, 8.9, 1.55, "Agents & skills\nacross the repositories", 32, WHITE, bold=True)
        self.text(
            slide,
            0.83,
            3.35,
            8.7,
            0.74,
            "How reusable engineering methods, governed orchestrators, and the neutral runtime fit together.",
            17,
            RGBColor(211, 225, 232),
        )
        self.pill(slide, 0.83, 4.62, 1.45, "CORE", CYAN)
        self.pill(slide, 2.48, 4.62, 1.65, "COMMUNITY", TEAL)
        self.pill(slide, 4.33, 4.62, 1.65, "ENTERPRISE", CORAL)
        self.text(
            slide,
            0.83,
            6.70,
            6.0,
            0.25,
            f"Workspace snapshot | {date.today().strftime('%d %b %Y')}",
            9,
            RGBColor(170, 194, 204),
        )

    def add_ecosystem_slide(self) -> None:
        """Add the repository-layer architecture slide."""
        slide = self.add_slide("One ecosystem, three content layers", "Architecture", 2)
        layers = [
            (
                "CORE",
                "neqsim",
                "Physics, canonical agents, mirrored skills, schemas, discovery and CI",
                CYAN,
                PALE_BLUE,
            ),
            (
                "COMMUNITY",
                "neqsim-community-agents\nneqsim-community-skills",
                "Public screening methods and transparent reusable workflows",
                TEAL,
                PALE_TEAL,
            ),
            (
                "ENTERPRISE",
                "neqsim-enterprise-agents\nneqsim-enterprise-skills",
                "Internal integrations, policy overlays and governed coordinators",
                CORAL,
                PALE_CORAL,
            ),
        ]
        for index, (label, repos, detail, accent, fill) in enumerate(layers):
            x = 0.72 + index * 4.18
            self.rect(slide, x, 1.72, 3.72, 3.18, fill, line_color=accent)
            self.pill(slide, x + 0.2, 1.92, 1.28, label, accent)
            self.text(slide, x + 0.2, 2.45, 3.25, 0.72, repos, 17, NAVY, bold=True)
            self.text(slide, x + 0.2, 3.30, 3.25, 0.84, detail, 12, INK)
            self.text(
                slide,
                x + 0.2,
                4.43,
                3.16,
                0.24,
                ["trust: core", "trust: community", "trust: internal"][index],
                9,
                accent,
                bold=True,
                font=MONO_FONT,
            )
            if index < 2:
                self.line(slide, x + 3.76, 3.28, x + 4.10, 3.28, MUTED, 1.8, arrow=True)
        self.rect(slide, 2.05, 5.39, 9.2, 1.0, WHITE, line_color=LINE)
        self.text(slide, 2.30, 5.60, 2.15, 0.25, "engineering-harness", 15, NAVY, bold=True)
        self.text(
            slide,
            4.47,
            5.54,
            6.45,
            0.48,
            "Neutral runtime: imports catalogs, validates permissions, executes workflows and writes evidence-backed reports.",
            12,
            INK,
        )

    def add_building_blocks_slide(self) -> None:
        """Add the conceptual building-block slide."""
        slide = self.add_slide("Each building block owns one kind of decision", "Operating model", 3)
        cards = [
            ("AGENT", "Who coordinates?", "Intent, scope, permissions, required skills, review gate", CYAN, PALE_BLUE),
            ("SKILL", "How is it done?", "Reusable engineering method, inputs, outputs and validation", TEAL, PALE_TEAL),
            ("TOOL", "What executes?", "NeqSim/MCP calculation, local package or governed data source", CORAL, PALE_CORAL),
            ("WORKFLOW / STUDY", "In what order?", "Steps, evidence, acceptance criteria and final assessment", YELLOW, PALE_YELLOW),
        ]
        for index, (label, question, answer, accent, fill) in enumerate(cards):
            x = 0.72 + index * 3.08
            self.rect(slide, x, 1.76, 2.74, 3.66, fill, line_color=accent)
            self.rect(slide, x + 0.18, 1.98, 0.54, 0.54, accent)
            self.text(slide, x + 0.18, 2.09, 0.54, 0.24, str(index + 1), 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
            self.text(slide, x + 0.18, 2.77, 2.30, 0.32, label, 13, accent, bold=True)
            self.text(slide, x + 0.18, 3.24, 2.30, 0.66, question, 18, NAVY, bold=True)
            self.text(slide, x + 0.18, 4.08, 2.30, 0.90, answer, 11.5, INK)
            if index < 3:
                self.line(slide, x + 2.77, 3.56, x + 3.02, 3.56, MUTED, 1.6, arrow=True)
        self.rect(slide, 1.62, 5.86, 10.1, 0.61, NAVY)
        self.text(
            slide,
            1.82,
            5.99,
            9.7,
            0.26,
            "Design rule: skills contain engineering methods; agents orchestrate skills.",
            14,
            WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def add_packaging_slide(self) -> None:
        """Add the repository packaging and schema slide."""
        slide = self.add_slide("Human-readable guidance, machine-readable contracts", "Packaging", 4)
        self.rect(slide, 0.72, 1.76, 5.78, 4.78, WHITE, line_color=LINE)
        self.text(slide, 0.98, 1.98, 5.1, 0.32, "Agent package", 18, NAVY, bold=True)
        tree = (
            "agents/<agent-name>/\n"
            "  AGENT.md          human workflow\n"
            "  agent.yaml        machine contract\n"
            "  prompts/          reusable prompts\n"
            "  workflows/        orchestration\n"
            "  examples/         worked cases\n"
            "  tests/            consistency checks"
        )
        self.text(slide, 1.02, 2.50, 5.0, 2.30, tree, 12.5, INK, font=MONO_FONT)
        self.pill(slide, 1.00, 5.55, 1.60, "required_skills", CYAN)
        self.pill(slide, 2.77, 5.55, 1.20, "extends", CORAL)
        self.pill(slide, 4.14, 5.55, 1.72, "human_review", TEAL)

        self.rect(slide, 6.82, 1.76, 5.79, 4.78, WHITE, line_color=LINE)
        self.text(slide, 7.08, 1.98, 5.1, 0.32, "Skill package", 18, NAVY, bold=True)
        skill_tree = (
            "skills/<category>/<skill>/\n"
            "  SKILL.md          method + triggers\n"
            "  package code      deterministic core\n"
            "  tests/            behavior checks\n"
            "  data/             safe fixtures"
        )
        self.text(slide, 7.12, 2.50, 4.95, 1.75, skill_tree, 12.5, INK, font=MONO_FONT)
        self.text(slide, 7.10, 4.56, 5.05, 0.32, "Catalog + schema layer", 14, NAVY, bold=True)
        self.bullet_list(
            slide,
            7.10,
            4.98,
            5.0,
            [
                "Root YAML catalogs enable discovery without scanning every package",
                "Canonical JSON schemas keep community and enterprise metadata aligned",
                "Namespaces separate public neqsim-* from internal enterprise-* skills",
            ],
            size=10.5,
            gap=0.42,
        )

    def add_snapshot_slide(self) -> None:
        """Add inventory metrics and category distribution."""
        slide = self.add_slide("A broad catalog, organized by trust and domain", "Workspace snapshot", 5)
        self.metric(slide, 0.72, 1.72, str(self.counts.total_agents), "agent package definitions", PALE_BLUE)
        self.metric(slide, 3.67, 1.72, str(self.counts.total_skills), "skill package definitions", PALE_TEAL)
        self.metric(slide, 6.62, 1.72, "3", "trust layers", PALE_CORAL)
        self.metric(slide, 9.57, 1.72, "1", "neutral runtime", PALE_YELLOW)

        self.text(slide, 0.72, 3.15, 5.60, 0.32, "Agent packages by repository", 16, NAVY, bold=True)
        agent_rows = [
            ("Core", self.counts.core_agents, CYAN),
            ("Community", self.counts.community_agents, TEAL),
            ("Enterprise", self.counts.enterprise_agents, CORAL),
        ]
        max_agent = max(value for _, value, _ in agent_rows) or 1
        for index, (label, value, color) in enumerate(agent_rows):
            y = 3.67 + index * 0.68
            self.text(slide, 0.72, y, 1.15, 0.25, label, 11, INK)
            self.rect(slide, 1.92, y + 0.03, 3.62 * value / max_agent, 0.23, color, radius=False)
            self.text(slide, 5.70, y - 0.01, 0.40, 0.25, str(value), 11, NAVY, bold=True, align=PP_ALIGN.RIGHT)

        self.text(slide, 6.72, 3.15, 5.52, 0.32, "External skills by category", 16, NAVY, bold=True)
        combined = self.counts.community_skill_categories + self.counts.enterprise_skill_categories
        category_rows = combined.most_common(6)
        max_category = max((value for _, value in category_rows), default=1)
        for index, (label, value) in enumerate(category_rows):
            y = 3.57 + index * 0.47
            color = TEAL if index % 2 == 0 else CYAN
            display = label.replace("-", " ").title()
            self.text(slide, 6.72, y, 1.62, 0.22, display, 9.5, INK)
            self.rect(slide, 8.36, y + 0.02, 3.26 * value / max_category, 0.18, color, radius=False)
            self.text(slide, 11.78, y - 0.02, 0.38, 0.22, str(value), 9.5, NAVY, bold=True, align=PP_ALIGN.RIGHT)
        self.text(
            slide,
            0.72,
            6.64,
            11.8,
            0.25,
            "Counts are checked-out package files as of build time; mirrored or shared capability IDs may overlap across repositories.",
            8.5,
            MUTED,
        )

    def add_extension_slide(self) -> None:
        """Add the enterprise extension-pattern slide."""
        slide = self.add_slide("Enterprise extends the public method with governance", "Overlay pattern", 6)
        self.rect(slide, 0.80, 1.82, 3.38, 3.85, PALE_TEAL, line_color=TEAL)
        self.pill(slide, 1.04, 2.06, 1.34, "COMMUNITY", TEAL)
        self.text(slide, 1.04, 2.62, 2.90, 0.66, "subsea-cooldown-agent", 17, NAVY, bold=True)
        self.text(slide, 1.04, 3.42, 2.84, 0.38, "required skill", 10, MUTED, bold=True)
        self.text(slide, 1.04, 3.83, 2.90, 0.72, "neqsim-surf-cooldown-screening", 12, TEAL, bold=True, font=MONO_FONT)
        self.text(slide, 1.04, 4.77, 2.90, 0.46, "Public physics + transparent screening", 11.5, INK)

        self.line(slide, 4.32, 3.74, 5.08, 3.74, CORAL, 2.4, arrow=True)
        self.text(slide, 4.34, 3.18, 0.72, 0.28, "extends", 10, CORAL, bold=True, align=PP_ALIGN.CENTER)

        self.rect(slide, 5.18, 1.82, 4.12, 3.85, PALE_CORAL, line_color=CORAL)
        self.pill(slide, 5.42, 2.06, 1.34, "ENTERPRISE", CORAL)
        self.text(slide, 5.42, 2.62, 3.54, 0.66, "surf-cooldown-agent", 17, NAVY, bold=True)
        self.text(slide, 5.42, 3.42, 3.34, 0.38, "effective skills = base + overlay", 10, MUTED, bold=True)
        self.text(slide, 5.42, 3.83, 3.45, 0.88, "neqsim-surf-cooldown-screening\nenterprise-surf-cooldown", 11.5, CORAL, bold=True, font=MONO_FONT)
        self.text(slide, 5.42, 4.89, 3.45, 0.36, "Policy verdict + review triggers", 11.5, INK)

        self.rect(slide, 9.72, 1.82, 2.74, 3.85, NAVY)
        self.text(slide, 9.98, 2.10, 2.16, 0.36, "Why this matters", 15, WHITE, bold=True)
        self.bullet_list(
            slide,
            9.98,
            2.76,
            2.15,
            [
                "One public method",
                "Additive policy layer",
                "Machine-checkable lineage",
                "Human review retained",
            ],
            size=11,
            color=WHITE,
            gap=0.60,
        )
        self.rect(slide, 1.76, 6.05, 9.85, 0.46, WHITE, line_color=LINE)
        self.text(
            slide,
            1.96,
            6.14,
            9.45,
            0.22,
            "Schema rule: an overlay may add required skills, but it cannot silently remove the base agent's dependencies.",
            10.5,
            NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def add_routing_slide(self) -> None:
        """Add the end-to-end discovery and execution flow."""
        slide = self.add_slide("From plain-language request to reviewed result", "Routing", 7)
        steps = [
            ("REQUEST", "Engineering question", CYAN, PALE_BLUE),
            ("DISCOVER", "router + semantic search", TEAL, PALE_TEAL),
            ("SELECT", "specialist agent", CORAL, PALE_CORAL),
            ("EXECUTE", "skills + tools", YELLOW, PALE_YELLOW),
            ("ASSESS", "evidence + human review", NAVY, WHITE),
        ]
        for index, (label, detail, accent, fill) in enumerate(steps):
            x = 0.58 + index * 2.55
            line_color = accent if accent != NAVY else NAVY
            self.rect(slide, x, 2.02, 2.18, 2.12, fill, line_color=line_color)
            self.text(slide, x + 0.16, 2.28, 1.86, 0.27, label, 10, accent if accent != NAVY else NAVY, bold=True)
            self.text(slide, x + 0.16, 2.80, 1.86, 0.86, detail, 15, NAVY, bold=True)
            if index < 4:
                self.line(slide, x + 2.20, 3.08, x + 2.50, 3.08, MUTED, 1.7, arrow=True)
        self.text(slide, 0.72, 4.73, 3.0, 0.28, "Example: pipeline + hydrate request", 13, NAVY, bold=True)
        self.rect(slide, 0.72, 5.18, 11.90, 1.12, WHITE, line_color=LINE)
        example_parts = [
            ("agent_search.py", CYAN),
            ("flow-assurance agent", TEAL),
            ("pressure-drop + hydrate skills", CORAL),
            ("NeqSim / MCP", YELLOW),
            ("review-ready assessment", NAVY),
        ]
        cursor = 0.98
        for index, (part, color) in enumerate(example_parts):
            width = [1.72, 2.18, 2.70, 1.68, 2.35][index]
            self.pill(slide, cursor, 5.57, width, part, color)
            cursor += width + 0.28
            if index < len(example_parts) - 1:
                self.text(slide, cursor - 0.22, 5.58, 0.20, 0.26, ">", 12, MUTED, bold=True, align=PP_ALIGN.CENTER)
        self.text(
            slide,
            0.72,
            6.62,
            11.7,
            0.26,
            "For multi-discipline work, compose a workflow in MCP or the engineering harness instead of forcing one agent to own every domain.",
            9.5,
            MUTED,
        )

    def add_governance_slide(self) -> None:
        """Add the governance and closing slide."""
        slide = self.add_slide("Validation keeps the catalog usable as it grows", "Governance", 8)
        checks = [
            ("SCHEMA", "Required metadata, namespaces, semver and review flags", CYAN),
            ("DEPENDENCIES", "Every required skill resolves across the combined index", TEAL),
            ("COOPERATION", "extends and coordinated_agents make relationships explicit", CORAL),
            ("DISCOVERY", "Search tests and generated maps keep packages findable", YELLOW),
            ("EVIDENCE", "Harness gates unsupported claims and preserves assumptions", NAVY),
        ]
        for index, (label, detail, color) in enumerate(checks):
            y = 1.77 + index * 0.86
            self.rect(slide, 0.78, y, 1.62, 0.52, color)
            self.text(slide, 0.78, y + 0.10, 1.62, 0.25, label, 9.5, WHITE, bold=True, align=PP_ALIGN.CENTER)
            self.text(slide, 2.67, y + 0.03, 5.05, 0.44, detail, 12, INK)
            if index < len(checks) - 1:
                self.line(slide, 1.59, y + 0.54, 1.59, y + 0.82, LINE, 1.3, arrow=True)

        self.rect(slide, 8.10, 1.77, 4.45, 4.83, NAVY)
        self.text(slide, 8.43, 2.12, 3.78, 0.40, "The organizing principle", 18, WHITE, bold=True)
        self.text(
            slide,
            8.43,
            2.88,
            3.68,
            1.46,
            "Put the method in a skill.\nPut coordination in an agent.\nPut execution and evidence in the harness.",
            18,
            WHITE,
            bold=True,
        )
        self.line(slide, 8.43, 4.64, 11.94, 4.64, RGBColor(82, 113, 128), 1.0)
        self.text(
            slide,
            8.43,
            4.96,
            3.63,
            1.06,
            "That separation makes capabilities reusable, enterprise policy additive, and review boundaries visible.",
            12.5,
            RGBColor(211, 225, 232),
        )
        self.text(
            slide,
            0.78,
            6.52,
            6.9,
            0.26,
            "Primary sources: repository catalogs, package manifests, canonical schemas and engineering-harness documentation.",
            8.5,
            MUTED,
        )

    def build(self) -> Presentation:
        """Build and return the complete presentation.

        Returns:
            Generated PowerPoint presentation.
        """
        self.add_title_slide()
        self.add_ecosystem_slide()
        self.add_building_blocks_slide()
        self.add_packaging_slide()
        self.add_snapshot_slide()
        self.add_extension_slide()
        self.add_routing_slide()
        self.add_governance_slide()
        return self.prs


def find_workspace_root(script_path: Path) -> Path:
    """Resolve the parent directory containing the sibling repositories.

    Args:
        script_path: Path to this generator script.

    Returns:
        Parent directory containing the core NeqSim repository.

    Raises:
        RuntimeError: If the expected workspace layout cannot be found.
    """
    for candidate in script_path.resolve().parents:
        if (candidate / "neqsim" / "pom.xml").exists():
            return candidate
    raise RuntimeError("Could not locate the workspace root containing neqsim/pom.xml")


def validate_inventory(counts: RepoCounts) -> None:
    """Reject an empty or partially resolved workspace inventory.

    Args:
        counts: Repository inventory to validate.

    Raises:
        RuntimeError: If a required repository package set is absent.
    """
    values = {
        "core agents": counts.core_agents,
        "core skills": counts.core_skills,
        "community agents": counts.community_agents,
        "community skills": counts.community_skills,
        "enterprise agents": counts.enterprise_agents,
        "enterprise skills": counts.enterprise_skills,
    }
    missing = [label for label, value in values.items() if value == 0]
    if missing:
        raise RuntimeError("Missing repository inventory: " + ", ".join(missing))


def build_deck(output_path: Path, workspace_root: Path) -> Path:
    """Generate and save the PowerPoint deck.

    Args:
        output_path: Destination PPTX file.
        workspace_root: Parent directory containing all source repositories.

    Returns:
        Written output path.
    """
    counts = RepoCounts(workspace_root)
    validate_inventory(counts)
    presentation = DeckBuilder(counts).build()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return output_path


def parse_args(arguments: Iterable[str] | None = None) -> argparse.Namespace:
    """Parse command-line arguments.

    Args:
        arguments: Optional explicit argument sequence.

    Returns:
        Parsed arguments.
    """
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, help="Destination .pptx path")
    parser.add_argument("--workspace-root", type=Path, help="Parent of the sibling repositories")
    return parser.parse_args(arguments)


def main() -> None:
    """Generate the presentation from the current workspace."""
    args = parse_args()
    script_path = Path(__file__)
    workspace_root = args.workspace_root or find_workspace_root(script_path)
    output_path = args.output or (
        workspace_root
        / "neqsim"
        / "docs"
        / "presentations"
        / "neqsim_agents_and_skills_organization.pptx"
    )
    written = build_deck(output_path.resolve(), workspace_root.resolve())
    print(str(written))


if __name__ == "__main__":
    main()
